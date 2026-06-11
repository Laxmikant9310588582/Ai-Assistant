import os
import uuid
import json
import pandas as pd
import pdfplumber
import httpx
from flask import Flask, request, jsonify, render_template, Blueprint, redirect, Response
from dotenv import load_dotenv
from werkzeug.serving import run_simple
from openai import OpenAI
import tiktoken
import fitz
from PIL import Image
from pptx import Presentation
import pytesseract
import psycopg2
import psycopg2.extras
from datetime import datetime

load_dotenv(override=True)

app = Flask(__name__)
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

client = OpenAI(
    api_key=os.getenv('OPENAI_API_KEY')
)

DB_CONFIG = {
    "host":     os.getenv("DB_HOST",     "localhost"),
    "database": os.getenv("DB_NAME",     "postgres"),
    "user":     os.getenv("DB_USER",     "postgres"),
    "password": os.getenv("DB_PASSWORD", "sharma123"),
    "port":     int(os.getenv("DB_PORT", 5432))
}

# .env mein ye lines add karo:
# DB_HOST=localhost
# DB_NAME=your_database_name
# DB_USER=your_db_user
# DB_PASSWORD=your_db_password
# DB_PORT=5432


# =====================================================
# DB CONNECTION
# =====================================================
def get_db_connection():
    try:
        conn = psycopg2.connect(
            host=os.getenv("DB_HOST", "localhost"),
            database=os.getenv("DB_NAME", "postgres"),
            user=os.getenv("DB_USER", "postgres"),
            password=os.getenv("DB_PASSWORD", ""),
            port=int(os.getenv("DB_PORT", 5432))
        )
        return conn
    except Exception as e:
        print(f"[DB] Connection error: {e}")
        return None


# =====================================================
# TABLE CREATION — startup pe ek baar chalega
# =====================================================
def create_table_if_not_exists():
    conn = get_db_connection()
    if not conn:
        print("[DB] ⚠️  Could not connect to DB — chat history will NOT be saved.")
        return
    try:
        cur = conn.cursor()
        cur.execute("""
            CREATE TABLE IF NOT EXISTS chat_history (
                id              SERIAL PRIMARY KEY,
                user_id         VARCHAR(255)    NOT NULL,
                project_code    VARCHAR(255)    DEFAULT '',
                feature         VARCHAR(50)     NOT NULL,
                filename        VARCHAR(500),
                user_message    TEXT,
                bot_response    TEXT            NOT NULL,
                created_at      TIMESTAMP       DEFAULT NOW()
            );

            CREATE INDEX IF NOT EXISTS idx_chat_history_user_project
                ON chat_history(user_id, project_code, created_at DESC);
        """)
        conn.commit()
        cur.close()
        conn.close()
        print("[DB] ✓ chat_history table ready.")
    except Exception as e:
        print(f"[DB] Table creation error: {e}")


# -------------------- Text Extraction --------------------

def clean_text(text):
    return ''.join(c if c.isprintable() else ' ' for c in text)


def ocr_image(image_path):
    img = Image.open(image_path)
    return clean_text(pytesseract.image_to_string(img))


def ocr_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = []
    for page in doc:
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        text.append(pytesseract.image_to_string(img))
    return "\n".join(text)


def load_data(file_path):
    if file_path.endswith('.txt'):
        return open(file_path, 'r', encoding='utf-8', errors='ignore').read()
    elif file_path.endswith(('.xlsx', '.xls')):
        return pd.read_excel(file_path).to_string(index=False)
    elif file_path.endswith('.pdf'):
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return clean_text(text) if text.strip() else ocr_pdf(file_path)
    elif file_path.endswith('.pptx'):
        prs = Presentation(file_path)
        return "\n".join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if hasattr(shape, "text")
        )
    elif file_path.endswith(('.png', '.jpg', '.jpeg')):
        return ocr_image(file_path)
    return ""


# -------------------- OpenAI Helpers --------------------

def split_text_into_token_chunks(text, max_tokens=3000, model="gpt-3.5-turbo"):
    enc = tiktoken.encoding_for_model(model)
    tokens = enc.encode(text)
    return [enc.decode(tokens[i:i + max_tokens]) for i in range(0, len(tokens), max_tokens)]


def query_openai_chunked(text, prompt):
    chunks = split_text_into_token_chunks(text)
    answers = []
    for chunk in chunks:
        messages = [
            {"role": "system", "content": "You are a helpful assistant. Use the provided context to answer the user's question accurately."},
            {"role": "system", "content": f"Context:\n{chunk}"},
            {"role": "user", "content": prompt}
        ]
        res = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=messages,
            max_tokens=500
        )
        answers.append(res.choices[0].message.content)
    return "\n\n".join(answers)


# -------------------- Metadata Helpers --------------------

def get_file_path(filename):
    meta_file = os.path.join(app.config['UPLOAD_FOLDER'], 'files_metadata.json')
    if not os.path.exists(meta_file):
        return None
    with open(meta_file) as f:
        for line in f:
            try:
                meta = json.loads(line)
                if meta['stored'] == filename:
                    return meta['file_path']
            except Exception:
                continue
    return None


# -------------------- Blueprint --------------------

ai = Blueprint('ai', __name__,
               static_folder='static',
               static_url_path='/static',
               template_folder='templates')


@ai.route('/')
def index():
    template_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'templates', 'index.html')
    with open(template_path, 'r', encoding='utf-8') as f:
        content = f.read()
    return Response(content, mimetype='text/html')


@ai.route('/upload', methods=['POST'])
def upload():
    try:
        file = request.files.get('file')
        if not file or not file.filename:
            return jsonify({'error': 'No file provided'}), 400

        ext = file.filename.rsplit('.', 1)[-1].lower()
        fname = f"{uuid.uuid4()}.{ext}"
        path = os.path.join(app.config['UPLOAD_FOLDER'], fname)
        file.save(path)

        metadata = {
            "original": file.filename,
            "stored": fname,
            "file_path": path
        }
        meta_file = os.path.join(app.config['UPLOAD_FOLDER'], 'files_metadata.json')
        with open(meta_file, 'a') as f:
            f.write(json.dumps(metadata) + "\n")

        return jsonify({
            "filename": fname,
            "original": file.filename,
            "success": True
        })

    except Exception as e:
        print(f"Upload error: {str(e)}")
        return jsonify({'error': f'Upload failed: {str(e)}'}), 500


@ai.route('/ask', methods=['POST'])
def ask():
    try:
        data = request.get_json()
        question = data.get('question')
        filename = data.get('filename')

        if not question or not filename:
            return jsonify({'error': 'Missing question or filename'}), 400

        file_path = get_file_path(filename)
        if not file_path or not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404

        text = load_data(file_path)
        if not text.strip():
            return jsonify({'error': 'Could not extract text from document'}), 400

        answer = query_openai_chunked(text, question)
        return jsonify({"answer": answer})

    except Exception as e:
        print(f"Ask error: {str(e)}")
        return jsonify({'error': f'Error processing question: {str(e)}'}), 500


@ai.route('/summary', methods=['POST'])
def summary():
    try:
        data = request.get_json()
        filename = data.get('filename')

        if not filename:
            return jsonify({'error': 'Missing filename'}), 400

        file_path = get_file_path(filename)
        if not file_path or not os.path.exists(file_path):
            return jsonify({'error': 'File not found'}), 404

        text = load_data(file_path)
        if not text.strip():
            return jsonify({'error': 'Could not extract text from document'}), 400

        prompt = (
            "Provide a clear structured summary with the following sections:\n"
            "1. High Level Overview\n"
            "2. Key Points\n"
            "3. Important Details"
        )
        summary_text = query_openai_chunked(text, prompt)
        return jsonify({"summary": summary_text})

    except Exception as e:
        print(f"Summary error: {str(e)}")
        return jsonify({'error': f'Error generating summary: {str(e)}'}), 500


# =====================================================
# ════════════════════════════════════════════════════
#  CHAT HISTORY ROUTES — NAYE ROUTES (ADDED)
# ════════════════════════════════════════════════════
# =====================================================

# ----------------------------------------------------------
#  POST /AI_Assistant/save-chat
#  Frontend se call hota hai jab bhi answer aata hai
#  Body: { user_id, project_code, feature,
#          filename, user_message, bot_response }
# ----------------------------------------------------------
@ai.route('/save-chat', methods=['POST'])
def save_chat():
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON"}), 400

    # Required fields check
    for field in ['user_id', 'feature', 'bot_response']:
        if not data.get(field):
            return jsonify({"error": f"Missing field: {field}"}), 400

    conn = get_db_connection()
    if not conn:
        return jsonify({"error": "Database connection failed"}), 500

    try:
        cur = conn.cursor()
        cur.execute("""
            INSERT INTO chat_history
                (user_id, project_code, feature, filename,
                 user_message, bot_response)
            VALUES (%s, %s, %s, %s, %s, %s)
            RETURNING id, created_at
        """, (
            data['user_id'],
            data.get('project_code', ''),
            data['feature'],           # 'qna' ya 'doc-summary'
            data.get('filename'),
            data.get('user_message'),  # QnA mein question, summary mein None
            data['bot_response']
        ))
        row = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()

        return jsonify({
            "success":    True,
            "id":         row[0],
            "created_at": str(row[1])
        }), 201

    except Exception as e:
        print(f"[save_chat] ERROR: {e}")
        conn.rollback()
        conn.close()
        return jsonify({"error": str(e)}), 500


# ----------------------------------------------------------
#  GET /AI_Assistant/recent-chats
#  Login ke baad frontend call karta hai
#  Returns: last 2 sessions with messages
# ----------------------------------------------------------
@ai.route('/recent-chats', methods=['GET'])
def recent_chats():
    user_id      = request.args.get('user_id', '').strip()
    project_code = request.args.get('project_code', '').strip()

    if not user_id:
        return jsonify({"error": "user_id required"}), 400

    conn = get_db_connection()
    if not conn:
        return jsonify({"error": "Database connection failed"}), 500

    try:
        cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)

        # Last 2 unique sessions (filename + feature + date)
        cur.execute("""
            SELECT DISTINCT
                filename,
                feature,
                DATE(created_at) AS chat_date
            FROM chat_history
            WHERE user_id      = %s
              AND project_code = %s
            ORDER BY chat_date DESC
            LIMIT 2
        """, (user_id, project_code))

        sessions = cur.fetchall()
        result   = []

        for session in sessions:
            cur.execute("""
                SELECT user_message, bot_response, created_at
                FROM chat_history
                WHERE user_id          = %s
                  AND project_code     = %s
                  AND filename         = %s
                  AND feature          = %s
                  AND DATE(created_at) = %s
                ORDER BY created_at ASC
            """, (
                user_id,
                project_code,
                session['filename'],
                session['feature'],
                session['chat_date']
            ))

            messages = cur.fetchall()
            result.append({
                "filename": session['filename'],
                "feature":  session['feature'],
                "date":     str(session['chat_date']),
                "messages": [
                    {
                        "user_message": m['user_message'],
                        "bot_response": m['bot_response'],
                        "time":         str(m['created_at'])
                    }
                    for m in messages
                ]
            })

        cur.close()
        conn.close()
        return jsonify(result), 200

    except Exception as e:
        print(f"[recent_chats] ERROR: {e}")
        conn.close()
        return jsonify({"error": str(e)}), 500


# ----------------------------------------------------------
#  GET /AI_Assistant/chat-history
#  History modal ke liye — saari history return karta hai
#  Query params: user_id, project_code, limit (default 200)
# ----------------------------------------------------------
@ai.route('/chat-history', methods=['GET'])
def chat_history():
    user_id      = request.args.get('user_id', '').strip()
    project_code = request.args.get('project_code', '').strip()
    limit        = min(int(request.args.get('limit', 200)), 500)

    if not user_id:
        return jsonify({"error": "user_id required"}), 400

    conn = get_db_connection()
    if not conn:
        return jsonify({"error": "Database connection failed"}), 500

    try:
        cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)
        cur.execute("""
            SELECT
                id, feature, filename,
                user_message, bot_response, created_at
            FROM chat_history
            WHERE user_id      = %s
              AND project_code = %s
            ORDER BY created_at DESC
            LIMIT %s
        """, (user_id, project_code, limit))

        rows = cur.fetchall()
        cur.close()
        conn.close()

        return jsonify([
            {
                "id":           r['id'],
                "feature":      r['feature'],
                "filename":     r['filename'],
                "user_message": r['user_message'],
                "bot_response": r['bot_response'],
                "created_at":   str(r['created_at'])
            }
            for r in rows
        ]), 200

    except Exception as e:
        print(f"[chat_history] ERROR: {e}")
        conn.close()
        return jsonify({"error": str(e)}), 500


# -------------------- App Setup --------------------

app.register_blueprint(ai, url_prefix='/AI_Assistant')


@app.route('/')
def root_redirect():
    return redirect('/AI_Assistant')

@app.route('/auth/logout')
def auth_logout():
    return redirect('/AI_Assistant/')


if __name__ == '__main__':
    # DB table startup pe bana do
    create_table_if_not_exists()

    print("\n" + "=" * 50)
    print("🤖 AI Assistant Started")
    print("=" * 50)
    print(f"📁 Upload folder: {UPLOAD_FOLDER}")
    print(f"🔑 OpenAI Key:  {'✓ Yes' if os.getenv('OPENAI_API_KEY') else '✗ NOT FOUND'}")
    print(f"🗄️  DB Host:     {os.getenv('DB_HOST', 'localhost')}")
    print(f"🗄️  DB Name:     {os.getenv('DB_NAME', 'NOT SET')}")
    print("🌐 URL: http://127.0.0.1:5009/AI_Assistant/")
    print("=" * 50 + "\n")

    run_simple('0.0.0.0', 5009, app, use_debugger=True, use_reloader=True)